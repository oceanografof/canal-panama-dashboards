"""
Aplicación Streamlit para Evaluación Hidrológica Modelo GLR.

Mejoras principales:
- Integra la hoja HostFolder con Resumen Simulación para recuperar
 hidrogeneración, vertidos y esclusajes reales.
- Separa completamente la información de Alhajuela y Gatún.
- Evita tablas excesivamente anchas y reemplaza valores faltantes por "—".
- Incluye gráficas legibles de nivel, caudales e hidrogeneración/vertidos.
- Exporta reportes legibles en PowerPoint, Word, PDF y HTML autocontenido.

Ejecución:
  streamlit run app_glr_mejorada.py
"""
from __future__ import annotations

import base64
import html
import io
import os
import re
import tempfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

# Exportadores opcionales: se validan al generar cada archivo.
try:
  from pptx import Presentation
  from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
  from pptx.util import Inches, Pt
  from pptx.dml.color import RGBColor
  from pptx.enum.shapes import MSO_SHAPE
except Exception: # pragma: no cover
  Presentation = None

try:
  from docx import Document
  from docx.enum.section import WD_ORIENT
  from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
  from docx.enum.text import WD_ALIGN_PARAGRAPH
  from docx.shared import Inches as DocxInches, Pt as DocxPt
  from docx.oxml import OxmlElement
  from docx.oxml.ns import qn
except Exception: # pragma: no cover
  Document = None

try:
  from reportlab.lib import colors
  from reportlab.lib.enums import TA_CENTER, TA_LEFT
  from reportlab.lib.pagesizes import A4, landscape
  from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
  from reportlab.lib.units import inch
  from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
  )
except Exception: # pragma: no cover
  SimpleDocTemplate = None


APP_TITLE = "Evaluación Hidrológica Modelo GLR / FEWS"
SYSTEM_LABEL = "Simulacion"
SUBTITLE = "Basado en pronóstico meteorológico - próximos 10 días"

# Paleta inspirada en la plantilla GLR, con mayor contraste.
BG = "#2f3135"
PANEL = "#3a3d42"
PANEL_2 = "#45494f"
BLUE = "#2d607c"
LIGHT_BLUE = "#8ec8ef"
GREEN = "#38a34a"
ORANGE = "#ff6b3a"
TAN = "#c5aa7a"
PURPLE = "#a98bd4"
RED = "#e35d6a"
WHITE = "#f7f7f7"
MUTED = "#c7cbd1"
GRID = "#697079"
DEFAULT_LOGO_NAME = "himh_logo.jpg"


@dataclass
class ReportMeta:
  titulo: str = APP_TITLE
  subtitulo: str = SUBTITLE
  hidrologo: str = "Jorge F. Rodriguez C."
  unidad: str = "HIMH - Sección de Hidrología"
  actualizacion: datetime = datetime.now()
  descripcion: str = "Resumen operativo de niveles, aportes, salidas, consumos, hidrogeneración y vertidos."
  horizonte_dias: int = 10
  logo_path: Optional[str] = None


def _default_logo_path() -> Optional[str]:
  """Busca un logo por defecto junto a la app para usarlo en los reportes."""
  base_dir = Path(__file__).resolve().parent
  candidates = [
    base_dir / DEFAULT_LOGO_NAME,
    base_dir / "HIMH_logo.jpg",
    base_dir / "HIMH_logo.png",
    base_dir / "logo_himh.jpg",
    base_dir / "logo_himh.png",
  ]
  for candidate in candidates:
    if candidate.exists():
      return str(candidate)
  return None


def resolve_logo_path(user_logo_path: Optional[str] = None) -> Optional[str]:
  if user_logo_path and os.path.exists(user_logo_path):
    return user_logo_path
  return _default_logo_path()


def _logo_data_uri(path: Optional[str]) -> str:
  if not path or not os.path.exists(path):
    return ""
  suffix = Path(path).suffix.lower()
  mime = "image/png" if suffix == ".png" else "image/jpeg"
  with open(path, "rb") as f:
    encoded = base64.b64encode(f.read()).decode("ascii")
  return f"data:{mime};base64,{encoded}"


# -----------------------------------------------------------------------------
# Lectura y normalización
# -----------------------------------------------------------------------------
def _clean_text(x) -> str:
  if pd.isna(x):
    return ""
  return str(x).replace("\n", " ").strip()


def _norm(s: str) -> str:
  s = str(s).lower()
  for a, b in [("ú", "u"), ("á", "a"), ("é", "e"), ("í", "i"), ("ó", "o"), ("ñ", "n")]:
    s = s.replace(a, b)
  s = re.sub(r"[^a-z0-9]+", " ", s)
  return re.sub(r"\s+", " ", s).strip()


def _first_col(cols: Iterable[str], patterns: Iterable[str]) -> Optional[str]:
  ncols = {c: _norm(c) for c in cols}
  for pat in patterns:
    p = _norm(pat)
    for c, nc in ncols.items():
      if p in nc:
        return c
  return None


def _to_numeric(s) -> pd.Series:
  return pd.to_numeric(s, errors="coerce")


def _read_sheet_names(excel_file) -> List[str]:
  return pd.ExcelFile(excel_file).sheet_names


def load_premisas(excel_file) -> List[str]:
  try:
    df = pd.read_excel(excel_file, sheet_name="Premisas")
  except Exception:
    return []
  premisas: List[str] = []
  for _, row in df.iterrows():
    values = [_clean_text(v) for v in row.tolist() if _clean_text(v)]
    txt = " ".join(values)
    txt = re.sub(r"^[•\-*\s]+", "", txt).strip()
    if not txt or txt.lower() == "premisas":
      continue
    ntxt = _norm(txt)
    # Excluir metadatos que algunas plantillas dejan debajo de las premisas.
    if "hidrologo de turno" in ntxt:
      continue
    if re.fullmatch(r"\d{4} \d{2} \d{2} \d{2} \d{2} \d{2}", ntxt):
      continue
    premisas.append(txt)
  return premisas[:12]


def _normalize_hostfolder(raw: pd.DataFrame) -> pd.DataFrame:
  cols = list(raw.columns)
  date_col = _first_col(cols, ["fecha hora", "fecha"]) or cols[0]
  mapping = {
    "fecha": date_col,
    "elev_alha_obs": _first_col(cols, ["elev observada madden", "elev observada alha"]),
    "elev_alha_proj": _first_col(cols, ["elev proyectada madden", "elev proyectada alha"]),
    "elev_gat_obs": _first_col(cols, ["elev observada gatun"]),
    "elev_gat_proj": _first_col(cols, ["elev proyectada gatun"]),
    "aportes_alha": _first_col(cols, ["aportes madden", "aportes alha"]),
    "aportes_gat": _first_col(cols, ["aportes gatun"]),
    "salidas_alha": _first_col(cols, ["salidas madden", "salidas alha"]),
    "salidas_gat": _first_col(cols, ["salidas gatun"]),
    "cg_alha": _first_col(cols, ["cg madden", "cg alha"]),
    "cg_gat_89": _first_col(cols, ["cg89 gatun", "cg 89 gatun"]),
    "cg_gat_885": _first_col(cols, ["cg88 5 gatun", "cg88.5 gatun"]),
    "obs_aportes_alha": _first_col(cols, ["observado aportes madden", "observado aportes alha"]),
    "obs_aportes_gat": _first_col(cols, ["observado aportes gatun"]),
    "trasvase_alha_gat": _first_col(cols, ["trasvase madden", "trasvase alha"]),
  }
  out = pd.DataFrame()
  out["fecha"] = pd.to_datetime(raw[mapping["fecha"]], errors="coerce")
  for key, col in mapping.items():
    if key == "fecha":
      continue
    out[key] = raw[col] if col in raw.columns else np.nan

  out["elev_alha"] = out["elev_alha_proj"].combine_first(out["elev_alha_obs"])
  out["elev_gat"] = out["elev_gat_proj"].combine_first(out["elev_gat_obs"])
  out["hidro_alha"] = np.nan
  out["hidro_gat"] = np.nan
  out["vert_alha"] = np.nan
  out["vert_gat"] = np.nan
  out["esclusajes"] = np.nan
  return out


def _normalize_resumen(raw: pd.DataFrame) -> pd.DataFrame:
  cols = list(raw.columns)
  c_fecha = "Fecha Hora " if "Fecha Hora " in cols else (cols[1] if len(cols) > 1 else cols[0])
  out = pd.DataFrame()
  out["fecha"] = pd.to_datetime(raw[c_fecha], errors="coerce")
  pos = {
    "elev_alha": 2,
    "elev_gat": 3,
    "aportes_alha": 4,
    "aportes_gat": 5,
    "hidro_alha": 6,
    "hidro_gat": 7,
    "vert_alha": 8,
    "vert_gat": 9,
    "esclusajes": 10,
  }
  for name, idx in pos.items():
    out[name] = raw.iloc[:, idx] if idx < raw.shape[1] else np.nan
  out["elev_alha_obs"] = out["elev_alha"]
  out["elev_alha_proj"] = out["elev_alha"]
  out["elev_gat_obs"] = out["elev_gat"]
  out["elev_gat_proj"] = out["elev_gat"]
  out["salidas_alha"] = np.nan
  out["salidas_gat"] = np.nan
  out["cg_alha"] = np.nan
  out["cg_gat_89"] = np.nan
  out["cg_gat_885"] = np.nan
  out["obs_aportes_alha"] = out["aportes_alha"]
  out["obs_aportes_gat"] = out["aportes_gat"]
  out["trasvase_alha_gat"] = np.nan
  return out


def _prepare_dataframe(df: pd.DataFrame) -> pd.DataFrame:
  if df.empty:
    return df
  df = df.dropna(subset=["fecha"]).sort_values("fecha").drop_duplicates("fecha", keep="last").reset_index(drop=True)
  df = df.loc[:, ~df.columns.duplicated()].copy()
  for c in df.columns:
    if c != "fecha":
      df[c] = _to_numeric(df[c])
  return df


def load_glr_excel(excel_file) -> Tuple[pd.DataFrame, List[str], str]:
  """Carga y combina HostFolder + Resumen Simulación cuando ambas existen."""
  sheet_names = _read_sheet_names(excel_file)
  host = pd.DataFrame()
  resumen = pd.DataFrame()

  if "HostFolder" in sheet_names:
    host = _prepare_dataframe(_normalize_hostfolder(pd.read_excel(excel_file, sheet_name="HostFolder")))
  if "Resumen Simulación" in sheet_names:
    resumen = _prepare_dataframe(_normalize_resumen(pd.read_excel(excel_file, sheet_name="Resumen Simulación")))

  if not host.empty and not resumen.empty:
    # HostFolder conserva niveles observados/proyectados, aportes y salidas totales.
    # Resumen Simulación aporta hidrogeneración, vertidos y esclusajes.
    ops = resumen[["fecha", "hidro_alha", "hidro_gat", "vert_alha", "vert_gat", "esclusajes"]].copy()
    df = host.drop(columns=["hidro_alha", "hidro_gat", "vert_alha", "vert_gat", "esclusajes"], errors="ignore")
    df = df.merge(ops, on="fecha", how="left", validate="one_to_one")
    source_sheet = "HostFolder + Resumen Simulación"
  elif not host.empty:
    df = host
    source_sheet = "HostFolder"
  elif not resumen.empty:
    df = resumen
    source_sheet = "Resumen Simulación"
  else:
    raise ValueError("No se encontró una hoja HostFolder o Resumen Simulación válida.")

  df = _prepare_dataframe(df)
  premisas = load_premisas(excel_file)
  return df, premisas, source_sheet


def filter_horizon(df: pd.DataFrame, start_date: Optional[pd.Timestamp], days: int) -> pd.DataFrame:
  if df.empty:
    return df.copy()
  start = pd.to_datetime(start_date if start_date is not None else df["fecha"].min()).normalize()
  end = start + pd.Timedelta(days=int(days))
  view = df[(df["fecha"] >= start) & (df["fecha"] < end)].copy()
  if view.empty:
    view = df.head(max(1, 24 * int(days))).copy()
  return view


# -----------------------------------------------------------------------------
# Resúmenes, métricas y tablas
# -----------------------------------------------------------------------------
def _safe_mean(df: pd.DataFrame, col: str) -> float:
  return float(df[col].mean()) if col in df and df[col].notna().any() else np.nan


def _safe_first(df: pd.DataFrame, col: str) -> float:
  s = df[col].dropna() if col in df else pd.Series(dtype=float)
  return float(s.iloc[0]) if not s.empty else np.nan


def _safe_last(df: pd.DataFrame, col: str) -> float:
  s = df[col].dropna() if col in df else pd.Series(dtype=float)
  return float(s.iloc[-1]) if not s.empty else np.nan


def report_metrics(df: pd.DataFrame) -> Dict[str, float]:
  if df.empty:
    return {}
  return {
    "alha_ini": _safe_first(df, "elev_alha"),
    "alha_fin": _safe_last(df, "elev_alha"),
    "gat_ini": _safe_first(df, "elev_gat"),
    "gat_fin": _safe_last(df, "elev_gat"),
    "ap_alha_prom": _safe_mean(df, "aportes_alha"),
    "ap_gat_prom": _safe_mean(df, "aportes_gat"),
    "sal_alha_prom": _safe_mean(df, "salidas_alha"),
    "sal_gat_prom": _safe_mean(df, "salidas_gat"),
    "consumo_prom": _safe_mean(df, "esclusajes"),
    "vol_esclusajes_dia": (_safe_mean(df, "esclusajes") * 86400.0 / 1_000_000.0) if pd.notna(_safe_mean(df, "esclusajes")) else np.nan,
    "vol_esclusajes_horizonte": (_safe_mean(df, "esclusajes") * float((df["fecha"].max() - df["fecha"].min()).total_seconds() + 3600.0) / 1_000_000.0) if ("fecha" in df and not df.empty and pd.notna(_safe_mean(df, "esclusajes"))) else np.nan,
    "hidro_alha_prom": _safe_mean(df, "hidro_alha"),
    "hidro_gat_prom": _safe_mean(df, "hidro_gat"),
    "vert_alha_prom": _safe_mean(df, "vert_alha"),
    "vert_gat_prom": _safe_mean(df, "vert_gat"),
  }


def _fmt(v: float, nd: int = 2, suffix: str = "") -> str:
  if v is None or pd.isna(v):
    return "N/D"
  return f"{v:.{nd}f}{suffix}"


def automatic_narrative(df: pd.DataFrame, premisas: List[str], meta: ReportMeta) -> str:
  """Interpreta el escenario sin repetir los valores mostrados en las tarjetas."""
  m = report_metrics(df)
  da = m.get("alha_fin", np.nan) - m.get("alha_ini", np.nan)
  dg = m.get("gat_fin", np.nan) - m.get("gat_ini", np.nan)

  def nivel_frase(nombre: str, variacion: float) -> str:
    if pd.isna(variacion):
      return f"Para {nombre} no hay información suficiente para establecer la tendencia del nivel"
    if abs(variacion) < 0.10:
      return f"{nombre} se mantiene prácticamente estable, con una variación de {variacion:+.2f} pies"
    if variacion > 0:
      return f"{nombre} presenta una tendencia ascendente, con una recuperación de {variacion:.2f} pies"
    return f"{nombre} presenta una tendencia descendente, con una reducción de {abs(variacion):.2f} pies"

  ap_alha = m.get("ap_alha_prom", np.nan)
  ap_gat = m.get("ap_gat_prom", np.nan)
  if pd.notna(ap_alha) and pd.notna(ap_gat):
    if abs(ap_alha - ap_gat) < 0.5:
      aporte_txt = "Los aportes previstos son similares en ambos embalses"
    elif ap_gat > ap_alha:
      aporte_txt = "Los aportes previstos son mayores en Gatún que en Alhajuela"
    else:
      aporte_txt = "Los aportes previstos son mayores en Alhajuela que en Gatún"
  else:
    aporte_txt = "No hay información completa para comparar los aportes de ambos embalses"

  hidro_alha = m.get("hidro_alha_prom", np.nan)
  hidro_gat = m.get("hidro_gat_prom", np.nan)
  if pd.notna(hidro_alha) and pd.notna(hidro_gat):
    if hidro_alha > 0.01 and hidro_gat <= 0.01:
      hidro_txt = "La hidrogeneración prevista se concentra en Alhajuela/Madden, sin generación en Gatún"
    elif hidro_gat > 0.01 and hidro_alha <= 0.01:
      hidro_txt = "La hidrogeneración prevista se concentra en Gatún, sin generación en Alhajuela/Madden"
    elif hidro_alha > 0.01 and hidro_gat > 0.01:
      hidro_txt = "Se prevé hidrogeneración en ambos embalses"
    else:
      hidro_txt = "No se prevé hidrogeneración en los embalses durante el horizonte evaluado"
  else:
    hidro_txt = "La información de hidrogeneración está incompleta para el horizonte evaluado"

  return (
    f"Para los próximos {meta.horizonte_dias} días, {nivel_frase('Alhajuela', da)}. "
    f"{nivel_frase('Gatún', dg)}. {aporte_txt}. {hidro_txt}."
  )


def _daily_base(df: pd.DataFrame) -> pd.DataFrame:
  if df.empty:
    return pd.DataFrame()
  work = df.copy().set_index("fecha")
  agg = {
    "elev_alha": "last",
    "elev_gat": "last",
    "aportes_alha": "mean",
    "aportes_gat": "mean",
    "hidro_alha": "mean",
    "hidro_gat": "mean",
    "esclusajes": "mean",
    "vert_gat": "mean",
    "vert_alha": "mean",
    "salidas_alha": "mean",
    "salidas_gat": "mean",
  }
  agg = {k: v for k, v in agg.items() if k in work.columns}
  return work.resample("D").agg(agg).dropna(how="all").reset_index()


def daily_summary_reservoir(df: pd.DataFrame, reservoir: str) -> pd.DataFrame:
  d = _daily_base(df)
  if d.empty:
    return d
  if reservoir.lower().startswith("gat"):
    cols = {
      "fecha": "Fecha",
      "elev_gat": "Nivel (pies)",
      "aportes_gat": "Aportes (m³/s)",
      "hidro_gat": "Hidro (MW)",
      "vert_gat": "Vertido (m³/s)",
      "esclusajes": "Esclusajes (m³/s)",
      "salidas_gat": "Salidas totales (m³/s)",
    }
  else:
    cols = {
      "fecha": "Fecha",
      "elev_alha": "Nivel (pies)",
      "aportes_alha": "Aportes (m³/s)",
      "hidro_alha": "Hidro (MW)",
      "vert_alha": "Vertido (m³/s)",
      "salidas_alha": "Salidas totales (m³/s)",
    }
  available = [c for c in cols if c in d.columns]
  out = d[available].rename(columns=cols)
  for c in out.columns:
    if c != "Fecha":
      out[c] = pd.to_numeric(out[c], errors="coerce").round(2)
  return out


def daily_summary_combined(df: pd.DataFrame) -> pd.DataFrame:
  d = _daily_base(df)
  if d.empty:
    return d
  cols = {
    "fecha": "Fecha",
    "elev_alha": "Alhajuela (pies)",
    "elev_gat": "Gatún (pies)",
    "aportes_alha": "AT. Alhajuela (m³/s)",
    "aportes_gat": "AT. Gatún (m³/s)",
    "hidro_alha": "Hidro Alhajuela (MW)",
    "hidro_gat": "Hidro Gatún (MW)",
    "vert_alha": "Vertido Alhajuela (m³/s)",
    "vert_gat": "Vertido Gatún (m³/s)",
    "esclusajes": "Esclusajes Gatún (m³/s)",
  }
  available = [c for c in cols if c in d.columns]
  out = d[available].rename(columns=cols)
  for c in out.columns:
    if c != "Fecha":
      out[c] = pd.to_numeric(out[c], errors="coerce").round(2)
  return out


def display_table(df: pd.DataFrame) -> pd.DataFrame:
  out = df.copy()
  if "Fecha" in out:
    out["Fecha"] = pd.to_datetime(out["Fecha"]).dt.strftime("%d/%m/%Y")
  for c in out.columns:
    if c != "Fecha":
      out[c] = out[c].map(lambda v: "—" if pd.isna(v) else f"{float(v):,.2f}")
  return out


# -----------------------------------------------------------------------------
# Gráficas Matplotlib para exportación
# -----------------------------------------------------------------------------
def _setup_dark_axis(ax, title: str, ylabel: str):
  ax.set_facecolor(BG)
  ax.figure.set_facecolor(BG)
  ax.set_title(title, color=WHITE, fontsize=16, fontweight="bold", loc="left", pad=11)
  ax.set_ylabel(ylabel, color=WHITE, fontsize=12.5, labelpad=11, fontweight="bold")
  ax.tick_params(axis="y", colors=MUTED, labelsize=11)
  ax.tick_params(axis="x", colors=MUTED, labelsize=11, pad=7)
  ax.grid(True, color=GRID, alpha=0.24, linestyle="--", linewidth=0.75)
  for sp in ax.spines.values():
    sp.set_color(GRID)


def _apply_time_axis(ax, x_values, labelsize: float = 10.0):
  x = pd.to_datetime(pd.Series(x_values)).dropna()
  if x.empty:
    return
  span_hours = max(1.0, (x.max() - x.min()).total_seconds() / 3600.0)
  if span_hours <= 48:
    locator = mdates.HourLocator(interval=3)
  elif span_hours <= 120:
    locator = mdates.HourLocator(interval=6)
  elif span_hours <= 240:
    locator = mdates.HourLocator(interval=12)
  else:
    locator = mdates.DayLocator(interval=1)
  ax.xaxis.set_major_locator(locator)
  ax.xaxis.set_major_formatter(mdates.DateFormatter("%d-%b\n%H:%M"))
  pad_hours = 2 if span_hours > 72 else 1
  ax.set_xlim(x.min() - pd.Timedelta(hours=pad_hours), x.max() + pd.Timedelta(hours=pad_hours))
  ax.tick_params(axis="x", colors=MUTED, labelsize=labelsize, pad=7)
  for label in ax.get_xticklabels():
    label.set_rotation(0)
    label.set_ha("center")


def _legend(ax, ncol: int = 3, y: float = 1.015):
  handles, labels = ax.get_legend_handles_labels()
  if not handles:
    return
  leg = ax.legend(
    handles,
    labels,
    loc="upper right",
    bbox_to_anchor=(1.0, y),
    ncol=min(ncol, len(handles)),
    frameon=False,
    fontsize=9.5,
    handlelength=2.4,
    columnspacing=1.35,
  )
  for item in leg.get_texts():
    item.set_color(WHITE)


def _padded_range(series_list, minimum_pad: float = 0.08, relative_pad: float = 0.12, floor_zero: bool = False):
  values = []
  for series in series_list:
    if series is None:
      continue
    s = pd.to_numeric(series, errors="coerce").dropna()
    if not s.empty:
      values.extend(s.tolist())
  if not values:
    return None
  lo, hi = float(min(values)), float(max(values))
  span = hi - lo
  pad = max(span * relative_pad, minimum_pad)
  if span == 0:
    pad = max(abs(hi) * 0.08, minimum_pad)
  lo2, hi2 = lo - pad, hi + pad
  if floor_zero:
    lo2 = min(0.0, lo2)
  return lo2, hi2


def _is_nearly_constant(series: pd.Series, tolerance: float = 1e-6) -> bool:
  s = pd.to_numeric(series, errors="coerce").dropna()
  return s.empty or float(s.max() - s.min()) <= tolerance


def plot_reservoir(df: pd.DataFrame, reservoir: str, output_path: str) -> str:
  """Gráfica exportable en dos paneles, similar a la plantilla GLR original."""
  is_gat = reservoir.lower().startswith("gat")
  if is_gat:
    elev, obs, aporte, salida, hydro, vert = "elev_gat", "elev_gat_obs", "aportes_gat", "salidas_gat", "hidro_gat", "vert_gat"
    title = "PROYECCIÓN HIDROLÓGICA EMBALSE GATÚN"
    flow_title = "Aportes, salidas totales y esclusajes"
  else:
    elev, obs, aporte, salida, hydro, vert = "elev_alha", "elev_alha_obs", "aportes_alha", "salidas_alha", "hidro_alha", "vert_alha"
    title = "PROYECCIÓN HIDROLÓGICA EMBALSE ALHAJUELA"
    flow_title = "Aportes y salidas Madden/Alhajuela"

  fig, (ax1, ax2) = plt.subplots(
    2, 1,
    figsize=(15.5, 8.25),
    sharex=True,
    facecolor=BG,
    gridspec_kw={"height_ratios": [1.02, 1.0], "hspace": 0.28},
  )
  fig.subplots_adjust(top=0.89, bottom=0.13, left=0.080, right=0.975)
  x = pd.to_datetime(df["fecha"])
  fig.suptitle(title, color=WHITE, fontsize=19, fontweight="bold", y=0.972)
  if not x.empty:
    fig.text(0.975, 0.942, f"Periodo: {x.min():%d/%m/%Y %H:%M} a {x.max():%d/%m/%Y %H:%M}",
             ha="right", va="center", color=MUTED, fontsize=9.2)

  # Panel 1: niveles, con límite ajustado a la serie principal.
  _setup_dark_axis(ax1, "Niveles proyectados y observados", "Nivel (pies)")
  main_level_series = []
  if elev in df and df[elev].notna().any():
    projected = pd.to_numeric(df[elev], errors="coerce")
    main_level_series.append(projected)
    ax1.plot(x, projected, color=GREEN, linewidth=2.7, label="Nivel proyectado", zorder=4)
    ax1.fill_between(x, projected, projected.min(), color=GREEN, alpha=0.10, zorder=1)
  if obs in df and df[obs].notna().any():
    observed = pd.to_numeric(df[obs], errors="coerce")
    main_level_series.append(observed)
    ax1.plot(x, observed, color=ORANGE, linewidth=2.3, label="Nivel observado", zorder=5)

  level_range = _padded_range(main_level_series, minimum_pad=0.10, relative_pad=0.16)
  if level_range:
    ax1.set_ylim(*level_range)
  ax1.yaxis.set_major_locator(mticker.MaxNLocator(nbins=6))
  ax1.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.2f"))
  # Resaltar el valor final proyectado sin saturar la gráfica.
  if elev in df and df[elev].notna().any():
    final_idx = pd.to_numeric(df[elev], errors="coerce").last_valid_index()
    if final_idx is not None:
      xf = pd.to_datetime(df.loc[final_idx, "fecha"])
      yf = float(df.loc[final_idx, elev])
      ax1.annotate(f"{yf:.2f}", xy=(xf, yf), xytext=(-5, 10), textcoords="offset points", ha="right", color=WHITE, fontsize=9.5, fontweight="bold")

  # Leyenda del nivel: únicamente datos observados y proyectados.
  h1, l1 = ax1.get_legend_handles_labels()
  if h1:
    leg = ax1.legend(h1, l1, loc="upper left", bbox_to_anchor=(0.0, 1.015), ncol=min(2, len(h1)), frameon=False, fontsize=10.5)
    for item in leg.get_texts():
      item.set_color(WHITE)

  # Panel 2: aportes y salidas, manteniendo todos los picos visibles.
  _setup_dark_axis(ax2, flow_title, "Caudal (m³/s)")
  flow_series = []
  if salida in df and df[salida].notna().any():
    sv = pd.to_numeric(df[salida], errors="coerce")
    flow_series.append(sv)
    ax2.plot(x, sv, color=BLUE, linewidth=2.6, label="Salidas totales", zorder=4)
    ax2.fill_between(x, 0, sv.fillna(0), color=BLUE, alpha=0.24, zorder=1)
  if aporte in df and df[aporte].notna().any():
    av = pd.to_numeric(df[aporte], errors="coerce")
    flow_series.append(av)
    ax2.plot(x, av, color=LIGHT_BLUE, linewidth=2.2, label="Aportes", zorder=5)
    ax2.fill_between(x, 0, av.fillna(0), color=LIGHT_BLUE, alpha=0.13, zorder=2)
  if is_gat and "esclusajes" in df and df["esclusajes"].notna().any():
    esc = pd.to_numeric(df["esclusajes"], errors="coerce")
    flow_series.append(esc)
    ax2.plot(x, esc, color=PURPLE, linewidth=2.1, linestyle="--", label="Esclusajes", zorder=6)

  flow_range = _padded_range(flow_series, minimum_pad=5.0, relative_pad=0.10, floor_zero=True)
  clipped_points = []
  if flow_range:
    all_flow = pd.concat([pd.to_numeric(s, errors="coerce").dropna() for s in flow_series], ignore_index=True) if flow_series else pd.Series(dtype=float)
    upper = flow_range[1]
    if not all_flow.empty:
      q995 = float(all_flow.quantile(0.995))
      vmax = float(all_flow.max())
      if q995 > 0 and vmax > q995 * 1.70:
        upper = q995 * 1.15
        # Marcar visualmente los picos extraordinarios que quedan fuera del rango.
        for source_name, series, color in [("Aportes", pd.to_numeric(df.get(aporte), errors="coerce"), LIGHT_BLUE), ("Salidas", pd.to_numeric(df.get(salida), errors="coerce"), BLUE)]:
          if series is None:
            continue
          mask = series > upper
          for xi, yi in zip(x[mask], series[mask]):
            clipped_points.append((xi, float(yi), source_name, color))
      ax2.set_ylim(0, upper)
    else:
      ax2.set_ylim(flow_range[0], flow_range[1])
  for xi, yi, source_name, color in clipped_points:
    top = ax2.get_ylim()[1]
    ax2.scatter([xi], [top * 0.965], marker="^", s=58, color=color, edgecolor=WHITE, linewidth=0.7, zorder=12, clip_on=False)
    ax2.annotate(f"{yi:,.0f}", xy=(xi, top * 0.965), xytext=(0, -18), textcoords="offset points", ha="center", va="top", color=WHITE, fontsize=8.8, fontweight="bold")
  if clipped_points:
    ax2.text(0.99, 0.03, "Triángulo = pico fuera de escala", transform=ax2.transAxes, ha="right", va="bottom", color=MUTED, fontsize=8.6)
  ax2.yaxis.set_major_locator(mticker.MaxNLocator(nbins=6))
  ax2.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.0f"))
  _legend(ax2, 3)

  _apply_time_axis(ax1, x, labelsize=9.2)
  _apply_time_axis(ax2, x, labelsize=10.0)
  ax1.tick_params(axis="x", labelbottom=True)
  ax2.set_xlabel("Fecha / hora", color=WHITE, fontsize=12.5, fontweight="bold", labelpad=10)
  fig.savefig(output_path, dpi=200, facecolor=BG, bbox_inches="tight")
  plt.close(fig)
  return output_path


def plot_operations_summary(df: pd.DataFrame, output_path: str) -> str:
  daily = _daily_base(df)
  fig, ax = plt.subplots(figsize=(10.5, 4.2), facecolor=BG)
  ax.set_facecolor(BG)
  if daily.empty:
    ax.text(0.5, 0.5, "Sin datos operativos", color=WHITE, ha="center", va="center")
  else:
    x = np.arange(len(daily))
    labels = pd.to_datetime(daily["fecha"]).dt.strftime("%d/%m")
    ha = pd.to_numeric(daily.get("hidro_alha", 0), errors="coerce").fillna(0)
    hg = pd.to_numeric(daily.get("hidro_gat", 0), errors="coerce").fillna(0)
    width = 0.34
    ax.bar(x - width / 2, ha, width, label="Hidro Alhajuela", color=TAN)
    ax.bar(x + width / 2, hg, width, label="Hidro Gatún", color=PURPLE)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=0, color=MUTED, fontsize=8.5)
    ax.set_ylabel("MW", color=WHITE)
    ax.tick_params(axis="y", colors=MUTED)
    ax.grid(axis="y", color=GRID, alpha=0.3, linestyle="--")
    ax2 = ax.twinx()
    esc = pd.to_numeric(daily.get("esclusajes", 0), errors="coerce").fillna(0)
    ax2.plot(x, esc, color=LIGHT_BLUE, marker="o", linewidth=2, label="Esclusajes")
    ax2.set_ylabel("Esclusajes (m³/s)", color=WHITE)
    ax2.tick_params(axis="y", colors=MUTED)
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    leg = ax.legend(h1 + h2, l1 + l2, loc="upper center", ncol=3, frameon=False, fontsize=8.5)
    for t in leg.get_texts():
      t.set_color(WHITE)
  ax.set_title("RESUMEN OPERATIVO DIARIO", color=WHITE, fontsize=14, fontweight="bold", pad=12)
  for sp in ax.spines.values():
    sp.set_color(GRID)
  fig.tight_layout()
  fig.savefig(output_path, dpi=180, facecolor=BG, bbox_inches="tight")
  plt.close(fig)
  return output_path


def make_chart_files(df: pd.DataFrame, out_dir: str) -> Dict[str, str]:
  Path(out_dir).mkdir(parents=True, exist_ok=True)
  return {
    "alha": plot_reservoir(df, "alhajuela", os.path.join(out_dir, "chart_alhajuela.png")),
    "gatun": plot_reservoir(df, "gatun", os.path.join(out_dir, "chart_gatun.png")),
    "ops": plot_operations_summary(df, os.path.join(out_dir, "chart_operaciones.png")),
  }


# -----------------------------------------------------------------------------
# Exportación PowerPoint
# -----------------------------------------------------------------------------
def _ppt_add_bg(slide):
  fill = slide.background.fill
  fill.solid()
  fill.fore_color.rgb = RGBColor(47, 49, 53)


def _ppt_text(slide, text, x, y, w, h, size=16, bold=False, color=(247, 247, 247), align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
  box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
  tf = box.text_frame
  tf.clear()
  tf.word_wrap = True
  tf.margin_left = Pt(3)
  tf.margin_right = Pt(3)
  tf.margin_top = Pt(2)
  tf.margin_bottom = Pt(2)
  tf.vertical_anchor = valign
  p = tf.paragraphs[0]
  p.alignment = align
  run = p.add_run()
  run.text = str(text)
  run.font.size = Pt(size)
  run.font.bold = bold
  run.font.color.rgb = RGBColor(*color)
  return box


def _ppt_header(slide, meta: ReportMeta, title: str):
  # Encabezado distribuido en zonas independientes para evitar traslapes.
  if meta.logo_path and os.path.exists(meta.logo_path):
    slide.shapes.add_picture(meta.logo_path, Inches(0.25), Inches(0.10), width=Inches(0.50))
  _ppt_text(slide, "CANAL DE PANAMÁ", 0.88, 0.14, 2.05, 0.30, 13.5, True)
  _ppt_text(slide, "HIMH - SECCIÓN DE HIDROLOGÍA", 0.88, 0.45, 2.40, 0.20, 8.2, True, (220, 224, 230))
  _ppt_text(slide, SYSTEM_LABEL, 5.35, 0.05, 2.65, 0.18, 7.8, True, (142, 200, 239), align=PP_ALIGN.CENTER)
  _ppt_text(slide, title, 3.25, 0.21, 6.85, 0.45, 16.0, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
  _ppt_text(slide, "Última actualización\n" + meta.actualizacion.strftime("%d/%m/%Y %H:%M"), 10.45, 0.12, 2.55, 0.52, 8.2, False, (215, 218, 224), align=PP_ALIGN.RIGHT)
  line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.76), Inches(12.78), Inches(0.025))
  line.fill.solid(); line.fill.fore_color.rgb = RGBColor(45, 96, 124); line.line.fill.background()


def _ppt_metric_card(slide, title: str, value: str, x: float, y: float, w: float = 2.25):
  shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.82))
  shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(58, 61, 66)
  shape.line.color.rgb = RGBColor(80, 86, 94)
  _ppt_text(slide, title, x + 0.12, y + 0.10, w - 0.24, 0.23, 8.5, False, (199, 203, 209))
  _ppt_text(slide, value, x + 0.12, y + 0.35, w - 0.24, 0.34, 16, True)


def _ppt_add_table(slide, df: pd.DataFrame, x: float, y: float, w: float, h: float, font_size: float = 8.2):
  shown = display_table(df)
  rows, cols = shown.shape[0] + 1, shown.shape[1]
  table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
  widths = [1.18] + [(w - 1.18) / max(cols - 1, 1)] * max(cols - 1, 0)
  for j, cw in enumerate(widths):
    table.columns[j].width = Inches(cw)
  values = [list(shown.columns)] + shown.astype(str).values.tolist()
  for i, row in enumerate(values):
    for j, value in enumerate(row):
      cell = table.cell(i, j)
      cell.text = str(value)
      cell.fill.solid()
      cell.fill.fore_color.rgb = RGBColor(45, 96, 124) if i == 0 else RGBColor(58, 61, 66) if i % 2 else RGBColor(69, 73, 79)
      cell.margin_left = Pt(2); cell.margin_right = Pt(2); cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
      cell.vertical_anchor = MSO_ANCHOR.MIDDLE
      for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
          run.font.size = Pt(font_size if i else font_size - 0.3)
          run.font.color.rgb = RGBColor(247, 247, 247)
          run.font.bold = i == 0
  return table


def generate_pptx(df: pd.DataFrame, premisas: List[str], meta: ReportMeta, output_path: str) -> str:
  if Presentation is None:
    raise RuntimeError("Falta python-pptx. Ejecute: py -m pip install python-pptx")
  with tempfile.TemporaryDirectory() as td:
    charts = make_chart_files(df, td)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    m = report_metrics(df)

    # 1. Portada / premisas / resumen
    slide = prs.slides.add_slide(blank); _ppt_add_bg(slide); _ppt_header(slide, meta, meta.titulo.upper())
    _ppt_text(slide, meta.subtitulo.upper(), 3.10, 0.95, 7.2, 0.36, 11.5, False, (220, 224, 230), PP_ALIGN.CENTER)
    level_cards = [
      ("Alhajuela inicial", _fmt(m.get("alha_ini"), 2, " pies")),
      ("Alhajuela final", _fmt(m.get("alha_fin"), 2, " pies")),
      ("Gatún inicial", _fmt(m.get("gat_ini"), 2, " pies")),
      ("Gatún final", _fmt(m.get("gat_fin"), 2, " pies")),
    ]
    for (lab, val), x in zip(level_cards, [0.65, 3.75, 6.85, 9.95]):
      _ppt_metric_card(slide, lab, val, x, 1.12, 2.75)

    flow_cards = [
      ("Aporte prom. Alhajuela", _fmt(m.get("ap_alha_prom"), 2, " m³/s")),
      ("Aporte prom. Gatún", _fmt(m.get("ap_gat_prom"), 2, " m³/s")),
      ("Volumen prom. esclusajes", _fmt(m.get("vol_esclusajes_dia"), 2, " hm³/día")),
    ]
    for (lab, val), x in zip(flow_cards, [1.85, 5.30, 8.75]):
      _ppt_metric_card(slide, lab, val, x, 2.06, 2.75)

    _ppt_text(slide, "PREMISAS OPERATIVAS", 0.70, 3.05, 4.0, 0.32, 13, True)
    items = premisas if premisas else ["Sin premisas cargadas desde la hoja Premisas."]
    y = 3.42
    for item in items[:6]:
      lines = max(1, int(np.ceil(len(item) / 95)))
      h = 0.24 + 0.14 * (lines - 1)
      shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(y), Inches(5.85), Inches(h + 0.06))
      shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(58, 61, 66); shape.line.color.rgb = RGBColor(70, 75, 82)
      _ppt_text(slide, "• " + item, 0.83, y + 0.02, 5.55, h, 8.8, False, valign=MSO_ANCHOR.MIDDLE)
      y += h + 0.11
      if y > 5.65:
        break

    _ppt_text(slide, "INTERPRETACIÓN DEL ESCENARIO", 6.90, 3.05, 4.8, 0.32, 13, True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(3.42), Inches(5.75), Inches(2.10))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(58, 61, 66); box.line.color.rgb = RGBColor(75, 80, 88)
    _ppt_text(slide, automatic_narrative(df, premisas, meta), 7.02, 3.60, 5.35, 1.72, 9.7, False)
    _ppt_text(slide, "Descripción del escenario", 6.90, 5.72, 3.2, 0.25, 10.5, True)
    _ppt_text(slide, meta.descripcion, 6.90, 5.98, 5.55, 0.48, 9.0, False, (225, 228, 233))
    _ppt_text(slide, "HIDRÓLOGO DE TURNO: " + meta.hidrologo, 3.8, 6.72, 5.7, 0.22, 8.6, False, (190, 196, 204), PP_ALIGN.CENTER)

    # 2 y 3. Gráficas por embalse
    for title, key in [("PROYECCIÓN HIDROLÓGICA - EMBALSE ALHAJUELA", "alha"), ("PROYECCIÓN HIDROLÓGICA - EMBALSE GATÚN", "gatun")]:
      slide = prs.slides.add_slide(blank); _ppt_add_bg(slide); _ppt_header(slide, meta, title)
      slide.shapes.add_picture(charts[key], Inches(0.38), Inches(0.96), width=Inches(12.55), height=Inches(6.22))

    # 4. Tabla consolidada de 10 días
    slide = prs.slides.add_slide(blank); _ppt_add_bg(slide); _ppt_header(slide, meta, "RESUMEN CONSOLIDADO - 10 DÍAS")
    _ppt_text(slide, f"Escenario: {meta.horizonte_dias} días - GLR", 0.55, 0.96, 3.0, 0.30, 10.2, True)
    _ppt_text(slide, f"HIDRÓLOGO DE TURNO: {meta.hidrologo}", 8.15, 0.96, 4.2, 0.30, 9.5, True, (220,224,230), PP_ALIGN.RIGHT)
    combined = daily_summary_combined(df).head(meta.horizonte_dias)
    _ppt_add_table(slide, combined, 0.30, 1.35, 12.72, 5.25, font_size=7.2)
    _ppt_text(slide, "Cuadro consolidado con niveles, aportes, hidrogeneración, vertidos y esclusajes para facilitar la lectura integral del escenario.",
         0.55, 6.75, 12.0, 0.22, 8.5, False, (195, 201, 208), PP_ALIGN.CENTER)

    prs.save(output_path)
  return output_path


# -----------------------------------------------------------------------------
# Exportación Word
# -----------------------------------------------------------------------------
def _docx_set_cell_shading(cell, fill: str):
  tc_pr = cell._tc.get_or_add_tcPr()
  shd = tc_pr.find(qn("w:shd"))
  if shd is None:
    shd = OxmlElement("w:shd")
    tc_pr.append(shd)
  shd.set(qn("w:fill"), fill)


def _docx_add_table(doc, df: pd.DataFrame):
  shown = display_table(df)
  table = doc.add_table(rows=1, cols=len(shown.columns))
  table.style = "Table Grid"
  table.autofit = True
  for j, c in enumerate(shown.columns):
    cell = table.rows[0].cells[j]
    cell.text = c
    _docx_set_cell_shading(cell, "2D607C")
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    for p in cell.paragraphs:
      p.alignment = WD_ALIGN_PARAGRAPH.CENTER
      for r in p.runs:
        r.font.bold = True; r.font.size = DocxPt(7.5); r.font.color.rgb = None
  for i, row in shown.iterrows():
    cells = table.add_row().cells
    for j, c in enumerate(shown.columns):
      cells[j].text = str(row[c])
      cells[j].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
      if i % 2:
        _docx_set_cell_shading(cells[j], "E9EDF1")
      for p in cells[j].paragraphs:
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in p.runs:
          r.font.size = DocxPt(7.2)
  return table


def generate_docx(df: pd.DataFrame, premisas: List[str], meta: ReportMeta, output_path: str) -> str:
  if Document is None:
    raise RuntimeError("Falta python-docx. Ejecute: py -m pip install python-docx")
  with tempfile.TemporaryDirectory() as td:
    charts = make_chart_files(df, td)
    doc = Document()
    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.LANDSCAPE
    sec.page_width, sec.page_height = sec.page_height, sec.page_width
    sec.top_margin = DocxInches(0.45); sec.bottom_margin = DocxInches(0.45)
    sec.left_margin = DocxInches(0.55); sec.right_margin = DocxInches(0.55)

    if meta.logo_path and os.path.exists(meta.logo_path):
      header = sec.header
      p_logo = header.paragraphs[0]
      p_logo.alignment = WD_ALIGN_PARAGRAPH.RIGHT
      run_logo = p_logo.add_run()
      run_logo.add_picture(meta.logo_path, width=DocxInches(0.60))

    m = report_metrics(df)
    sys_p = doc.add_paragraph(SYSTEM_LABEL)
    sys_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in sys_p.runs:
      run.bold = True
      run.font.size = DocxPt(10)
    title = doc.add_heading(meta.titulo, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = doc.add_paragraph(meta.subtitulo); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Última actualización: {meta.actualizacion:%d/%m/%Y %H:%M} | Hidrólogo de turno: {meta.hidrologo}")
    tmet = doc.add_table(rows=2, cols=7)
    tmet.style = "Table Grid"
    headers = ["Alhajuela inicial", "Alhajuela final", "Gatún inicial", "Gatún final", "Aporte prom. Alhajuela", "Aporte prom. Gatún", "Vol. esclusajes"]
    vals = [_fmt(m.get("alha_ini"),2," pies"), _fmt(m.get("alha_fin"),2," pies"), _fmt(m.get("gat_ini"),2," pies"), _fmt(m.get("gat_fin"),2," pies"), _fmt(m.get("ap_alha_prom"),2," m³/s"), _fmt(m.get("ap_gat_prom"),2," m³/s"), _fmt(m.get("vol_esclusajes_dia"),2," hm³/día")]
    for i, h in enumerate(headers): tmet.cell(0,i).text = h
    for i, v in enumerate(vals): tmet.cell(1,i).text = v
    doc.add_heading("Premisas operativas", level=1)
    for item in premisas if premisas else ["Sin premisas cargadas."]:
      doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Interpretación del escenario", level=1)
    doc.add_paragraph(automatic_narrative(df, premisas, meta))

    doc.add_page_break()
    doc.add_heading("Proyección hidrológica - Embalse Alhajuela", level=1)
    doc.add_picture(charts["alha"], width=DocxInches(10.0))

    doc.add_page_break()
    doc.add_heading("Proyección hidrológica - Embalse Gatún", level=1)
    doc.add_picture(charts["gatun"], width=DocxInches(10.0))

    doc.add_page_break()
    doc.add_heading("Resumen consolidado de 10 días", level=1)
    doc.add_paragraph(f"Escenario: {meta.horizonte_dias} días - GLR")
    _docx_add_table(doc, daily_summary_combined(df).head(meta.horizonte_dias))
    doc.add_paragraph(meta.descripcion)
    doc.save(output_path)
  return output_path


# -----------------------------------------------------------------------------
# Exportación PDF con flujo automático (sin superposición)
# -----------------------------------------------------------------------------
def generate_pdf(df: pd.DataFrame, premisas: List[str], meta: ReportMeta, output_path: str) -> str:
  if SimpleDocTemplate is None:
    raise RuntimeError("Falta reportlab. Ejecute: py -m pip install reportlab")
  with tempfile.TemporaryDirectory() as td:
    charts = make_chart_files(df, td)
    page_w, page_h = landscape(A4)
    doc = SimpleDocTemplate(
      output_path,
      pagesize=landscape(A4),
      leftMargin=0.42 * inch,
      rightMargin=0.42 * inch,
      topMargin=0.86 * inch,
      bottomMargin=0.42 * inch,
      title=meta.titulo,
      author=meta.hidrologo,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleWhite", parent=styles["Title"], textColor=colors.white, fontName="Helvetica-Bold", fontSize=18, leading=22, alignment=TA_CENTER, spaceAfter=12)
    h1 = ParagraphStyle("H1White", parent=styles["Heading1"], textColor=colors.white, fontName="Helvetica-Bold", fontSize=14, leading=17, spaceAfter=8)
    body = ParagraphStyle("BodyWhite", parent=styles["BodyText"], textColor=colors.HexColor(WHITE), fontSize=9.2, leading=13, spaceAfter=6)
    bullet = ParagraphStyle("BulletWhite", parent=body, leftIndent=14, firstLineIndent=-8, bulletIndent=2)
    small = ParagraphStyle("SmallWhite", parent=body, fontSize=7.6, leading=9.2)

    def page_bg(canvas_obj, doc_obj):
      canvas_obj.saveState()
      canvas_obj.setFillColor(colors.HexColor(BG))
      canvas_obj.rect(0, 0, page_w, page_h, fill=1, stroke=0)
      canvas_obj.setStrokeColor(colors.HexColor(BLUE))
      canvas_obj.setLineWidth(1.2)
      canvas_obj.line(0.42 * inch, page_h - 0.70 * inch, page_w - 0.42 * inch, page_h - 0.70 * inch)
      canvas_obj.setFillColor(colors.white)
      canvas_obj.setFont("Helvetica-Bold", 10)
      canvas_obj.drawString(0.92 * inch, page_h - 0.48 * inch, "CANAL DE PANAMÁ | HIMH - SECCIÓN DE HIDROLOGÍA")
      canvas_obj.setFont("Helvetica-Bold", 8)
      canvas_obj.setFillColor(colors.HexColor(LIGHT_BLUE))
      canvas_obj.drawCentredString(page_w / 2.0, page_h - 0.48 * inch, SYSTEM_LABEL)
      if meta.logo_path and os.path.exists(meta.logo_path):
        try:
          canvas_obj.drawImage(meta.logo_path, 0.45 * inch, page_h - 0.69 * inch, width=0.34 * inch, height=0.34 * inch, preserveAspectRatio=True, mask='auto')
        except Exception:
          pass
      canvas_obj.setFont("Helvetica", 7.5)
      canvas_obj.setFillColor(colors.HexColor(MUTED))
      canvas_obj.drawRightString(page_w - 0.45 * inch, page_h - 0.48 * inch, meta.actualizacion.strftime("Actualización: %d/%m/%Y %H:%M"))
      canvas_obj.drawRightString(page_w - 0.45 * inch, 0.22 * inch, f"Página {doc_obj.page}")
      canvas_obj.restoreState()

    m = report_metrics(df)
    story = [Paragraph(SYSTEM_LABEL, ParagraphStyle("SystemLabel", parent=body, textColor=colors.HexColor(LIGHT_BLUE), fontName="Helvetica-Bold", fontSize=9, alignment=TA_CENTER, spaceAfter=3)), Paragraph(meta.titulo, title_style), Paragraph(meta.subtitulo, body), Spacer(1, 0.06 * inch)]
    metric_table = Table([
      ["Alhajuela inicial", "Alhajuela final", "Gatún inicial", "Gatún final", "Aporte prom. Alhajuela", "Aporte prom. Gatún", "Vol. esclusajes"],
      [_fmt(m.get("alha_ini"),2," pies"), _fmt(m.get("alha_fin"),2," pies"), _fmt(m.get("gat_ini"),2," pies"), _fmt(m.get("gat_fin"),2," pies"), _fmt(m.get("ap_alha_prom"),2," m³/s"), _fmt(m.get("ap_gat_prom"),2," m³/s"), _fmt(m.get("vol_esclusajes_dia"),2," hm³/día")],
    ], colWidths=[1.47*inch]*7)
    metric_table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor(BLUE)),("BACKGROUND",(0,1),(-1,1),colors.HexColor(PANEL)),("TEXTCOLOR",(0,0),(-1,-1),colors.white),("GRID",(0,0),(-1,-1),0.35,colors.HexColor('#70757D')),("ALIGN",(0,0),(-1,-1),'CENTER')]))
    story += [metric_table, Spacer(1,0.10*inch), Paragraph("Premisas operativas", h1)]
    for item in premisas if premisas else ["Sin premisas cargadas desde la hoja Premisas."]:
      story.append(Paragraph("• " + html.escape(item), bullet))
    story += [Spacer(1,0.08*inch), Paragraph("Interpretación del escenario", h1), Paragraph(html.escape(automatic_narrative(df, premisas, meta)), body)]

    story.extend([PageBreak(), Paragraph("Proyección hidrológica - Embalse Alhajuela", h1), RLImage(charts["alha"], width=10.65 * inch, height=6.15 * inch)])
    story.extend([PageBreak(), Paragraph("Proyección hidrológica - Embalse Gatún", h1), RLImage(charts["gatun"], width=10.65 * inch, height=6.15 * inch)])
    story.extend([PageBreak(), Paragraph("Resumen consolidado de 10 días", h1), Paragraph(html.escape(f"Escenario: {meta.horizonte_dias} días - GLR"), body)])
    story.append(_pdf_table(daily_summary_combined(df).head(meta.horizonte_dias), small))
    story.append(Spacer(1,0.08*inch))
    story.append(Paragraph(html.escape(meta.descripcion), body))
    doc.build(story, onFirstPage=page_bg, onLaterPages=page_bg)
  return output_path


def _pdf_table(df: pd.DataFrame, style_paragraph: ParagraphStyle):
  shown = display_table(df)
  data = [[Paragraph(html.escape(str(c)).replace(" ", "<br/>", 1), style_paragraph) for c in shown.columns]]
  for _, row in shown.iterrows():
    data.append([Paragraph(html.escape(str(row[c])), style_paragraph) for c in shown.columns])
  n = len(shown.columns)
  first = 0.92 * inch
  remaining = 10.55 * inch - first
  col_widths = [first] + [remaining / max(n - 1, 1)] * max(n - 1, 0)
  table = Table(data, colWidths=col_widths, repeatRows=1, hAlign="CENTER")
  table.setStyle(TableStyle([
    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BLUE)),
    ("TEXTCOLOR", (0, 0), (-1, -1), colors.white),
    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
    ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor(PANEL)),
    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor(PANEL), colors.HexColor(PANEL_2)]),
    ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#70757D")),
    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ("TOPPADDING", (0, 0), (-1, -1), 5),
    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
  ]))
  return table


# -----------------------------------------------------------------------------
# Exportación HTML autocontenida
# -----------------------------------------------------------------------------
def _img_data_uri(path: str) -> str:
  mime = "image/png"
  with open(path, "rb") as f:
    encoded = base64.b64encode(f.read()).decode("ascii")
  return f"data:{mime};base64,{encoded}"


def _html_table(df: pd.DataFrame) -> str:
  shown = display_table(df)
  return shown.to_html(index=False, escape=True, border=0, classes="data-table")


def generate_html(df: pd.DataFrame, premisas: List[str], meta: ReportMeta, output_path: str) -> str:
  with tempfile.TemporaryDirectory() as td:
    charts = make_chart_files(df, td)
    m = report_metrics(df)
    logo_tag = (f'<img class="logo-mini" src="{_logo_data_uri(meta.logo_path)}" alt="Logo HIMH">' if meta.logo_path and os.path.exists(meta.logo_path) else "")
    cards = [
      ("Alhajuela inicial", _fmt(m.get("alha_ini"), 2, " pies")),
      ("Alhajuela final", _fmt(m.get("alha_fin"), 2, " pies")),
      ("Gatún inicial", _fmt(m.get("gat_ini"), 2, " pies")),
      ("Gatún final", _fmt(m.get("gat_fin"), 2, " pies")),
      ("Aporte promedio de Alhajuela", _fmt(m.get("ap_alha_prom"), 2, " m³/s")),
      ("Aporte promedio de Gatún", _fmt(m.get("ap_gat_prom"), 2, " m³/s")),
      ("Volumen promedio de esclusajes", _fmt(m.get("vol_esclusajes_dia"), 2, " hm³/día")),
    ]
    cards_html = "".join(f'<div class="card"><span>{html.escape(k)}</span><strong>{html.escape(v)}</strong></div>' for k, v in cards)
    prem_html = "".join(f"<li>{html.escape(p)}</li>" for p in (premisas if premisas else ["Sin premisas cargadas."]))
    table_html = _html_table(daily_summary_combined(df).head(meta.horizonte_dias))
    doc_html = f"""<!doctype html>
<html lang="es">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{html.escape(meta.titulo)}</title>
<style>
:root{{--bg:{BG};--panel:{PANEL};--panel2:{PANEL_2};--blue:{BLUE};--text:{WHITE};--muted:{MUTED};}}
*{{box-sizing:border-box}} body{{margin:0;background:var(--bg);color:var(--text);font-family:Arial,Helvetica,sans-serif;line-height:1.42}}
.container{{max-width:1450px;margin:auto;padding:24px}} header{{border-bottom:3px solid var(--blue);padding-bottom:18px;margin-bottom:24px}}
.header-top{{display:flex;align-items:flex-start;justify-content:space-between;gap:18px}} .header-copy{{flex:1}} .logo-mini{{width:62px;max-width:62px;height:auto;opacity:.96;flex:0 0 auto}}
.system-label{{text-align:center;color:var(--light-blue,#8ec8ef);font-weight:700;letter-spacing:1.2px;font-size:14px;margin-bottom:4px}}
h1{{font-size:30px;margin:0 0 6px;text-align:center}} .subtitle{{text-align:center;color:var(--muted);margin:0}}
.meta{{display:flex;justify-content:space-between;gap:12px;flex-wrap:wrap;margin-top:16px;color:var(--muted);font-size:13px}}
.cards{{display:grid;grid-template-columns:repeat(auto-fit,minmax(190px,1fr));gap:12px;margin:22px 0}}
.card{{background:var(--panel);border:1px solid #555b63;border-radius:10px;padding:14px 16px;min-height:78px}} .card span{{display:block;color:var(--muted);font-size:13px;margin-bottom:5px}} .card strong{{font-size:21px}}
section{{background:var(--panel);border:1px solid #50555d;border-radius:12px;padding:20px;margin:18px 0;break-inside:avoid;page-break-inside:avoid}}
h2{{margin:0 0 14px;font-size:22px;border-left:5px solid var(--blue);padding-left:10px}} ul{{padding-left:24px}} li{{margin:7px 0}} .narrative{{font-size:16px;background:var(--panel2);padding:16px;border-radius:8px}}
.chart{{width:100%;height:auto;display:block;border-radius:8px;background:#2f3135}} .table-wrap{{overflow-x:auto}} .data-table{{width:100%;border-collapse:collapse;font-size:12px;white-space:nowrap}}
.data-table th{{background:var(--blue);padding:8px 6px;text-align:center}} .data-table td{{padding:7px 6px;text-align:right;border-bottom:1px solid #5a6068}} .data-table td:first-child{{text-align:center}} .data-table tr:nth-child(even){{background:var(--panel2)}}
button{{position:fixed;right:18px;bottom:18px;background:var(--blue);color:white;border:0;border-radius:8px;padding:12px 16px;font-weight:bold;cursor:pointer;box-shadow:0 4px 12px #0008}}
.page{{page-break-after:always}} .page:last-of-type{{page-break-after:auto}} footer{{text-align:center;color:var(--muted);font-size:12px;padding:18px}}
@media(max-width:900px){{.container{{padding:12px}} h1{{font-size:24px}} .header-top{{gap:10px}} .logo-mini{{width:46px;max-width:46px}}}}
@media print{{button{{display:none}} body{{background:white;color:#111}} section,.card{{background:white;color:#111;border-color:#bbb}} .subtitle,.meta,.card span,footer{{color:#444}} .narrative{{background:#f2f2f2}} .data-table td{{border-color:#ccc}} .chart{{page-break-inside:avoid}}}}
</style>
</head>
<body>
<div class="container">
<div class="page">
<header><div class="header-top"><div class="header-copy"><div class="system-label">{SYSTEM_LABEL}</div><h1>{html.escape(meta.titulo)}</h1><p class="subtitle">{html.escape(meta.subtitulo)}</p><div class="meta"><span>Hidrólogo de turno: {html.escape(meta.hidrologo)}</span><span>Actualización: {meta.actualizacion:%d/%m/%Y %H:%M}</span></div></div>{logo_tag}</div></header>
<div class="cards">{cards_html}</div>
<section><h2>Interpretación del escenario</h2><div class="narrative">{html.escape(automatic_narrative(df, premisas, meta))}</div></section>
<section><h2>Premisas operativas</h2><ul>{prem_html}</ul><p>{html.escape(meta.descripcion)}</p></section>
</div>
<div class="page"><section><h2>Embalse Alhajuela</h2><img class="chart" src="{_img_data_uri(charts['alha'])}" alt="Gráficas de Alhajuela"></section></div>
<div class="page"><section><h2>Embalse Gatún</h2><img class="chart" src="{_img_data_uri(charts['gatun'])}" alt="Gráficas de Gatún"></section></div>
<div class="page"><section><h2>Resumen consolidado de 10 días</h2><p><strong>Escenario:</strong> {meta.horizonte_dias} días - GLR</p><div class="table-wrap">{table_html}</div></section><footer>Actualización: {meta.actualizacion:%d/%m/%Y %H:%M}</footer></div>
</div><button onclick="window.print()">Imprimir / Guardar PDF</button>
</body></html>"""
    Path(output_path).write_text(doc_html, encoding="utf-8")
  return output_path


def _bytes_from_file(path: str) -> bytes:
  with open(path, "rb") as f:
    return f.read()


# -----------------------------------------------------------------------------
# Interfaz Streamlit
# -----------------------------------------------------------------------------
def run_streamlit():
  import streamlit as st
  try:
    import plotly.graph_objects as go
    from plotly.subplots import make_subplots
  except Exception:
    go = None
    make_subplots = None

  st.set_page_config(page_title="GLR Reportes", layout="wide", initial_sidebar_state="expanded")
  st.markdown("""
  <style>
   .block-container {padding-top: 1.35rem; padding-bottom: 2rem; max-width: 1550px;}
   div[data-testid="stMetric"] {background:#f7f8fa; border:1px solid #e1e4e8; padding:12px 14px; border-radius:9px;}
   div[data-testid="stMetric"] label {font-size:0.82rem;}
   .stTabs [data-baseweb="tab-list"] {gap:8px; flex-wrap:wrap;}
   .stTabs [data-baseweb="tab"] {height:42px; padding:0 16px; white-space:nowrap;}
   div[data-testid="stDataFrame"] {border:1px solid #d9dde3; border-radius:8px; overflow:hidden;}
   h2, h3 {letter-spacing:-0.2px;}
  </style>
  """, unsafe_allow_html=True)

  st.markdown(f"<div style='text-align:center;color:{LIGHT_BLUE};font-weight:700;letter-spacing:1.2px'>{SYSTEM_LABEL}</div>", unsafe_allow_html=True)
  st.title(APP_TITLE)
  st.caption("Dashboard integral separado por embalse y exportación a PowerPoint, Word, PDF y HTML")

  with st.sidebar:
    st.header("1. Archivo fuente")
    uploaded = st.file_uploader("Cargar Modelo_GLR.xlsx", type=["xlsx", "xlsm", "xls"])
    st.header("2. Configuración")
    hidrologo = st.text_input("Hidrólogo de turno", value="Jorge F. Rodriguez C.")
    horizonte = st.number_input("Días a resumir", min_value=1, max_value=30, value=10, step=1)
    descripcion = st.text_area("Descripción breve del escenario", value="10 días con datos de CHPS-FEWS")
    logo = st.file_uploader("Logo opcional PNG/JPG", type=["png", "jpg", "jpeg"])

  if not uploaded:
    st.info("Cargue el Excel del Modelo GLR para generar el dashboard y los reportes.")
    st.stop()

  try:
    df, premisas, source_sheet = load_glr_excel(uploaded)
  except Exception as exc:
    st.error(f"No se pudo leer el Excel: {exc}")
    st.stop()
  if df.empty:
    st.error("El archivo no contiene datos válidos de fecha o niveles.")
    st.stop()

  min_date = df["fecha"].min().date(); max_date = df["fecha"].max().date()
  with st.sidebar:
    start_date = st.date_input("Fecha inicial del reporte", value=min_date, min_value=min_date, max_value=max_date)

  view = filter_horizon(df, pd.Timestamp(start_date), int(horizonte))
  user_logo_path = None
  if logo is not None:
    suffix = Path(logo.name).suffix or ".png"
    tmp_logo = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    tmp_logo.write(logo.read()); tmp_logo.close(); user_logo_path = tmp_logo.name
  logo_path = resolve_logo_path(user_logo_path)

  meta = ReportMeta(
    hidrologo=hidrologo,
    actualizacion=datetime.now(),
    descripcion=descripcion,
    horizonte_dias=int(horizonte),
    logo_path=logo_path,
  )
  m = report_metrics(view)
  st.success(f"Datos integrados desde: {source_sheet}. Registros del horizonte: {len(view):,}")
  if logo_path and os.path.exists(logo_path):
    st.caption("Logo HIMH usado en los reportes")
    st.image(logo_path, width=78)
  if not view.get("hidro_alha", pd.Series(dtype=float)).notna().any():
    st.warning("No se encontró una serie de hidrogeneración para Alhajuela en el periodo seleccionado.")

  tabs = st.tabs(["Resumen", "Embalse Alhajuela", "Embalse Gatún", "Premisas", "Cuadro 10 días", "Exportar"])

  with tabs[0]:
    st.subheader("Indicadores principales")
    da = m.get("alha_fin", np.nan) - m.get("alha_ini", np.nan)
    dg = m.get("gat_fin", np.nan) - m.get("gat_ini", np.nan)

    n1, n2, n3, n4 = st.columns(4)
    n1.metric("Alhajuela inicial", _fmt(m.get("alha_ini"), 2, " pies"))
    n2.metric("Alhajuela final", _fmt(m.get("alha_fin"), 2, " pies"), _fmt(da, 2, " pies"))
    n3.metric("Gatún inicial", _fmt(m.get("gat_ini"), 2, " pies"))
    n4.metric("Gatún final", _fmt(m.get("gat_fin"), 2, " pies"), _fmt(dg, 2, " pies"))

    a1, a2, a3 = st.columns(3)
    a1.metric("Aporte promedio de Alhajuela", _fmt(m.get("ap_alha_prom"), 2, " m³/s"))
    a2.metric("Aporte promedio de Gatún", _fmt(m.get("ap_gat_prom"), 2, " m³/s"))
    a3.metric("Volumen promedio de esclusajes", _fmt(m.get("vol_esclusajes_dia"), 2, " hm³/día"))

    st.subheader("Interpretación del escenario")
    st.info(automatic_narrative(view, premisas, meta))

    st.subheader("Operación prevista: hidrogeneración y vertidos")
    d = _daily_base(view)
    if not d.empty:
      op1, op2, op3 = st.columns(3)
      op1.metric("Hidrogeneración promedio en Alhajuela", _fmt(pd.to_numeric(d.get("hidro_alha"), errors="coerce").mean(), 2, " MW"))
      op2.metric("Hidrogeneración promedio en Gatún", _fmt(pd.to_numeric(d.get("hidro_gat"), errors="coerce").mean(), 2, " MW"))
      op3.metric("Vertido promedio total", _fmt((pd.to_numeric(d.get("vert_alha"), errors="coerce").fillna(0) + pd.to_numeric(d.get("vert_gat"), errors="coerce").fillna(0)).mean(), 2, " m³/s"))

      op_table = pd.DataFrame({
        "Fecha": pd.to_datetime(d["fecha"]),
        "Hidro Alhajuela (MW)": pd.to_numeric(d.get("hidro_alha"), errors="coerce"),
        "Hidro Gatún (MW)": pd.to_numeric(d.get("hidro_gat"), errors="coerce"),
        "Esclusajes (m³/s)": pd.to_numeric(d.get("esclusajes"), errors="coerce"),
        "Vertido Alhajuela (m³/s)": pd.to_numeric(d.get("vert_alha"), errors="coerce"),
        "Vertido Gatún (m³/s)": pd.to_numeric(d.get("vert_gat"), errors="coerce"),
      })
      st.dataframe(display_table(op_table.head(int(horizonte))), use_container_width=True, hide_index=True, height=355)

      constant_notes = []
      for label, col, unit in [
        ("hidrogeneración de Alhajuela", "hidro_alha", "MW"),
        ("hidrogeneración de Gatún", "hidro_gat", "MW"),
        ("esclusajes", "esclusajes", "m³/s"),
      ]:
        s = pd.to_numeric(d.get(col, pd.Series(dtype=float)), errors="coerce").dropna()
        if not s.empty and _is_nearly_constant(s):
          constant_notes.append(f"{label}: {s.iloc[0]:.2f} {unit}, constante durante el horizonte")
      if constant_notes:
        st.info(". ".join(constant_notes) + ".")

  def reservoir_tab(reservoir: str):
    is_gat = reservoir == "Gatún"
    prefix = "gat" if is_gat else "alha"
    elev = f"elev_{prefix}"; obs = f"elev_{prefix}_obs"; aporte = f"aportes_{prefix}"; salida = f"salidas_{prefix}"
    hydro = f"hidro_{prefix}"; vert = f"vert_{prefix}"
    d = _daily_base(view)
    rm = report_metrics(view)

    st.subheader(f"Embalse {reservoir}")
    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Nivel inicial", _fmt(rm.get(f"{prefix}_ini"), 2, " pies"))
    k2.metric("Nivel final", _fmt(rm.get(f"{prefix}_fin"), 2, " pies"))
    k3.metric("Aporte promedio", _fmt(rm.get(f"ap_{prefix}_prom"), 2, " m³/s"))
    k4.metric("Hidrogeneración promedio", _fmt(rm.get(f"hidro_{prefix}_prom"), 2, " MW"))

    if go:
      st.markdown("#### Nivel")
      fig = go.Figure()
      fig.add_trace(go.Scatter(x=view["fecha"], y=view[elev], mode="lines", name="Nivel proyectado", line=dict(color=GREEN, width=3)))
      if obs in view and view[obs].notna().any():
        fig.add_trace(go.Scatter(x=view["fecha"], y=view[obs], mode="lines", name="Nivel observado", line=dict(color=ORANGE, width=2)))
      level_values = [pd.to_numeric(view[elev], errors="coerce")]
      if obs in view and view[obs].notna().any():
        level_values.append(pd.to_numeric(view[obs], errors="coerce"))
      level_rng = _padded_range(level_values, minimum_pad=0.08, relative_pad=0.16)
      fig.update_layout(
        height=430,
        margin=dict(l=65, r=35, t=45, b=65),
        yaxis_title="Nivel (pies)",
        legend=dict(orientation="h", y=1.15, font=dict(size=12)),
        hovermode="x unified",
        font=dict(size=13),
        xaxis=dict(tickformat="%d-%b", dtick=86400000, tickangle=0, tickfont=dict(size=12), automargin=True, hoverformat="%d/%m/%Y %H:%M"),
      )
      if level_rng:
        fig.update_yaxes(range=list(level_rng), tickformat=".2f", tickfont=dict(size=12), automargin=True)
      st.plotly_chart(fig, use_container_width=True)

      st.markdown("#### Aportes")
      aporte_values = pd.to_numeric(view[aporte], errors="coerce")
      a1, a2, a3 = st.columns(3)
      a1.metric("Aporte promedio", _fmt(aporte_values.mean(), 2, " m³/s"))
      a2.metric("Aporte máximo", _fmt(aporte_values.max(), 2, " m³/s"))
      a3.metric("Aporte mínimo", _fmt(aporte_values.min(), 2, " m³/s"))
      fig2 = go.Figure()
      fig2.add_trace(go.Scatter(x=view["fecha"], y=aporte_values, name="Aportes", mode="lines", fill="tozeroy", line=dict(color=LIGHT_BLUE, width=2.7)))
      aporte_rng = _padded_range([aporte_values], minimum_pad=5.0, relative_pad=0.12, floor_zero=True)
      fig2.update_layout(
        height=390,
        margin=dict(l=65, r=35, t=35, b=65),
        yaxis_title="Aportes (m³/s)",
        legend=dict(orientation="h", y=1.14, font=dict(size=12)),
        hovermode="x unified",
        font=dict(size=13),
        xaxis=dict(tickformat="%d-%b", dtick=86400000, tickangle=0, tickfont=dict(size=12), automargin=True, hoverformat="%d/%m/%Y %H:%M"),
      )
      if aporte_rng:
        fig2.update_yaxes(range=[max(0, aporte_rng[0]), aporte_rng[1]], tickformat=".0f", tickfont=dict(size=12), automargin=True)
      else:
        fig2.update_yaxes(rangemode="tozero", tickfont=dict(size=12), automargin=True)
      st.plotly_chart(fig2, use_container_width=True)

      st.markdown("#### Salidas y consumo")
      fig_out = go.Figure()
      if salida in view and view[salida].notna().any():
        fig_out.add_trace(go.Scatter(x=view["fecha"], y=view[salida], name="Salidas totales", mode="lines", line=dict(color=BLUE, width=3)))
      if is_gat and view.get("esclusajes", pd.Series(dtype=float)).notna().any():
        fig_out.add_trace(go.Scatter(x=view["fecha"], y=view["esclusajes"], name="Esclusajes", mode="lines", line=dict(color=PURPLE, width=2, dash="dash")))
      fig_out.update_layout(
        height=390, margin=dict(l=65, r=35, t=35, b=65), yaxis_title="Caudal (m³/s)",
        legend=dict(orientation="h", y=1.14, font=dict(size=12)), hovermode="x unified", font=dict(size=13),
        xaxis=dict(tickformat="%d-%b", dtick=86400000, tickangle=0, tickfont=dict(size=12), automargin=True, hoverformat="%d/%m/%Y %H:%M"),
      )
      fig_out.update_yaxes(rangemode="tozero", tickfont=dict(size=12), automargin=True)
      st.plotly_chart(fig_out, use_container_width=True)

      st.markdown("#### Hidrogeneración y vertidos")
      if not d.empty:
        hs = pd.to_numeric(d.get(hydro, pd.Series(dtype=float)), errors="coerce")
        vs = pd.to_numeric(d.get(vert, pd.Series(dtype=float)), errors="coerce")
        hcol, vcol = st.columns(2)
        hcol.metric("Hidrogeneración promedio", _fmt(hs.mean(), 2, " MW"), help="Promedio diario del horizonte seleccionado")
        vcol.metric("Vertido promedio", _fmt(vs.mean(), 2, " m³/s"), help="Promedio diario del horizonte seleccionado")
        hv_table = pd.DataFrame({"Fecha": pd.to_datetime(d["fecha"]), "Hidrogeneración (MW)": hs, "Vertido (m³/s)": vs})
        st.dataframe(display_table(hv_table.head(int(horizonte))), use_container_width=True, hide_index=True, height=300)
        if (not _is_nearly_constant(hs)) or (not _is_nearly_constant(vs)):
          fig3 = make_subplots(specs=[[{"secondary_y": True}]])
          fig3.add_trace(go.Bar(x=d["fecha"], y=hs, name="Hidrogeneración (MW)", marker_color=TAN), secondary_y=False)
          fig3.add_trace(go.Scatter(x=d["fecha"], y=vs, name="Vertido (m³/s)", mode="lines+markers", line=dict(color=RED, width=2.5)), secondary_y=True)
          fig3.update_layout(height=380, margin=dict(l=55, r=55, t=35, b=40), legend=dict(orientation="h", y=1.14), hovermode="x unified")
          fig3.update_yaxes(title_text="MW", secondary_y=False, rangemode="tozero")
          fig3.update_yaxes(title_text="m³/s", secondary_y=True, rangemode="tozero")
          st.plotly_chart(fig3, use_container_width=True)
        else:
          st.caption("Las variables se mantienen constantes; la tabla comunica mejor el escenario que una gráfica sin variación.")
    else:
      st.line_chart(view.set_index("fecha")[[elev, aporte, salida]])

    st.markdown(f"#### Tabla resumida - {reservoir}")
    st.dataframe(display_table(daily_summary_reservoir(view, reservoir).head(int(horizonte))), use_container_width=True, hide_index=True, height=410)

  with tabs[1]:
    reservoir_tab("Alhajuela")
  with tabs[2]:
    reservoir_tab("Gatún")

  with tabs[3]:
    st.subheader("Premisas operativas")
    if premisas:
      for i, item in enumerate(premisas, 1):
        st.markdown(f"**{i}.** {item}")
    else:
      st.warning("No se encontraron premisas en la hoja Premisas.")

  with tabs[4]:
    st.subheader("Cuadro consolidado de 10 días")
    st.dataframe(display_table(daily_summary_combined(view).head(int(horizonte))), use_container_width=True, hide_index=True, height=470)
    st.caption("Se presenta un solo cuadro consolidado, similar al formato operativo GLR, con niveles, aportes, hidrogeneración, vertidos y esclusajes.")

  with tabs[5]:
    st.subheader("Exportar reporte")
    st.write("Los reportes se generan en una estructura compacta de 4 hojas: 1) portada/resumen y premisas, 2) gráfico de Alhajuela, 3) gráfico de Gatún y 4) cuadro consolidado de 10 días.")
    if st.button("Generar archivos de reporte", type="primary", use_container_width=True):
      with st.spinner("Generando PowerPoint, Word, PDF y HTML..."):
        exports = {}
        errors = {}
        stamp = meta.actualizacion.strftime("%Y%m%d_%H%M")
        export_names = {
          "pptx": f"Reporte_GLR_{stamp}.pptx",
          "docx": f"Reporte_GLR_{stamp}.docx",
          "pdf": f"Reporte_GLR_{stamp}.pdf",
          "html": f"Reporte_GLR_{stamp}.html",
        }
        with tempfile.TemporaryDirectory() as td:
          targets = {
            "pptx": (generate_pptx, os.path.join(td, export_names["pptx"])),
            "docx": (generate_docx, os.path.join(td, export_names["docx"])),
            "pdf": (generate_pdf, os.path.join(td, export_names["pdf"])),
            "html": (generate_html, os.path.join(td, export_names["html"])),
          }
          for ext, (func, path) in targets.items():
            try:
              func(view, premisas, meta, path)
              exports[ext] = _bytes_from_file(path)
            except Exception as exc:
              errors[ext] = str(exc)
        st.session_state["glr_exports"] = exports
        st.session_state["glr_export_errors"] = errors
        st.session_state["glr_export_names"] = export_names

    exports = st.session_state.get("glr_exports", {})
    errors = st.session_state.get("glr_export_errors", {})
    export_names = st.session_state.get("glr_export_names", {})
    if exports:
      cols = st.columns(4)
      if "pptx" in exports:
        cols[0].download_button("Descargar PowerPoint", exports["pptx"], export_names.get("pptx", "Reporte_GLR.pptx"), mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
      if "docx" in exports:
        cols[1].download_button("Descargar Word", exports["docx"], export_names.get("docx", "Reporte_GLR.docx"), mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
      if "pdf" in exports:
        cols[2].download_button("Descargar PDF", exports["pdf"], export_names.get("pdf", "Reporte_GLR.pdf"), mime="application/pdf", use_container_width=True)
      if "html" in exports:
        cols[3].download_button("Descargar HTML", exports["html"], export_names.get("html", "Reporte_GLR.html"), mime="text/html", use_container_width=True)
    for ext, err in errors.items():
      st.error(f"No se pudo generar {ext.upper()}: {err}")


if __name__ == "__main__":
  import sys
  if len(sys.argv) >= 3 and sys.argv[1] == "--export-demo":
    excel = sys.argv[2]
    out_dir = sys.argv[3] if len(sys.argv) > 3 else "."
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    data, prem, source = load_glr_excel(excel)
    view = filter_horizon(data, None, 10)
    meta = ReportMeta(hidrologo="Jorge F. Rodriguez C.", horizonte_dias=10, logo_path=resolve_logo_path())
    print(f"Fuente: {source}; registros: {len(view)}")
    stamp = meta.actualizacion.strftime("%Y%m%d_%H%M")
    print(generate_pptx(view, prem, meta, os.path.join(out_dir, f"Reporte_GLR_{stamp}.pptx")))
    print(generate_docx(view, prem, meta, os.path.join(out_dir, f"Reporte_GLR_{stamp}.docx")))
    print(generate_pdf(view, prem, meta, os.path.join(out_dir, f"Reporte_GLR_{stamp}.pdf")))
    print(generate_html(view, prem, meta, os.path.join(out_dir, f"Reporte_GLR_{stamp}.html")))
  else:
    run_streamlit()
